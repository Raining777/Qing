"""Universal file parser: PDF, DOCX, PPTX, images (Vision), code, text."""
import base64
import hashlib
import logging
from pathlib import Path
from typing import Optional

from app.config import CODE_EXTENSIONS, MAX_FILE_SIZE_MB

logger = logging.getLogger(__name__)


def get_file_type(file_path: str) -> str:
    """Detect file type. Returns: text_pdf | scanned_pdf | docx | pptx | image | code | text | unsupported"""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return "text_pdf"  # scanned detection happens during parse
    if ext == ".docx":
        return "docx"
    if ext == ".pptx":
        return "pptx"
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        return "image"
    if ext in CODE_EXTENSIONS:
        return "code"
    if ext in {".txt", ".md"}:
        return "text"
    return "unsupported"


def file_md5(file_path: str) -> str:
    """Compute MD5 hash for deduplication."""
    h = hashlib.md5()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def file_size_ok(file_path: str) -> bool:
    """Check if file is within size limit."""
    mb = Path(file_path).stat().st_size / (1024 * 1024)
    return mb <= MAX_FILE_SIZE_MB


# ── PDF ──

def parse_pdf(file_path: str) -> tuple[str, bool]:
    """Extract text from PDF. Returns (text, is_scanned).

    If PyMuPDF returns very little text, marks as scanned.
    """
    import fitz
    text_parts = []
    total_chars = 0
    try:
        doc = fitz.open(file_path)
        for page in doc:
            page_text = page.get_text()
            text_parts.append(page_text)
            total_chars += len(page_text.strip())
        doc.close()
    except Exception as e:
        logger.error(f"PDF parse error: {e}")
        return "", True

    full_text = "\n\n".join(text_parts)
    # If average < 50 chars per page, likely scanned
    pages = len(text_parts) if text_parts else 1
    is_scanned = (total_chars / max(pages, 1)) < 50
    return full_text, is_scanned


def parse_pdf_first_pages(file_path: str, pages: int = 2) -> str:
    """Extract first N pages of PDF for classification."""
    import fitz
    try:
        doc = fitz.open(file_path)
        text_parts = []
        for i, page in enumerate(doc):
            if i >= pages:
                break
            text_parts.append(page.get_text())
        doc.close()
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"PDF first-pages error: {e}")
        return ""


async def parse_scanned_pdf(file_path: str, llm) -> str:
    """Extract text from scanned PDF using Claude Vision.

    Processes pages in batches of 5 to minimize API calls.
    """
    import fitz
    from app.services.image_utils import compress_image

    try:
        doc = fitz.open(file_path)
        all_text = []
        batch_size = 5

        for batch_start in range(0, len(doc), batch_size):
            batch_pages = []
            for i in range(batch_start, min(batch_start + batch_size, len(doc))):
                page = doc[i]
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("jpeg")
                img_bytes = compress_image(img_bytes, max_px=1200)
                b64 = base64.b64encode(img_bytes).decode()
                batch_pages.append(b64)

            # Send batch to Vision
            prompt = (
                "请从这些教材或学习资料页面中提取全部可见文字。"
                "保留标题、段落、列表、公式和原有结构；输出纯文本，方便后续总结和复习。"
            )
            for j, b64 in enumerate(batch_pages):
                try:
                    page_text = await llm.chat_with_image(
                        system="你负责准确提取学习资料页面中的文字和结构，不要加入原图中没有的内容。",
                        messages=[{"role": "user", "content": prompt}],
                        image_b64=b64,
                        mime="image/jpeg",
                    )
                    all_text.append(page_text)
                except Exception as e:
                    logger.error(f"Scanned page {batch_start + j + 1} error: {e}")
                    all_text.append(f"[Page {batch_start + j + 1} extraction failed]")

        doc.close()
        return "\n\n".join(all_text)
    except Exception as e:
        logger.error(f"Scanned PDF error: {e}")
        return ""


# ── DOCX ──

def parse_docx(file_path: str) -> str:
    """Extract text from Word document."""
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.error(f"DOCX parse error: {e}")
        return ""


def parse_docx_first_pages(file_path: str, paragraphs: int = 20) -> str:
    """Extract first N paragraphs for classification."""
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n\n".join(p.text for p in list(doc.paragraphs)[:paragraphs] if p.text.strip())
    except Exception:
        return ""


# ── PPTX ──

def parse_pptx(file_path: str) -> str:
    """Extract text from PowerPoint."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text.append(text)
            if slide_text:
                parts.append(f"## Slide {slide_num}\n" + "\n".join(slide_text))
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"PPTX parse error: {e}")
        return ""


def parse_pptx_first_pages(file_path: str, slides: int = 3) -> str:
    """Extract first N slides for classification."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            if slide_num > slides:
                break
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t:
                            slide_text.append(t)
            if slide_text:
                parts.append("\n".join(slide_text))
        return "\n\n".join(parts)
    except Exception:
        return ""


# ── Image (via Vision) ──

async def parse_image(file_path: str, llm) -> str:
    """Extract content from image using Claude Vision."""
    from app.services.image_utils import load_and_compress_b64
    try:
        b64, mime = load_and_compress_b64(file_path)
        prompt = (
            "请彻底分析这张学习资料图片，并提取成适合总结和复习的结构化文本：\n"
            "1. 所有可见文字\n"
            "2. 图表、流程图、示意图或曲线图的结构与含义\n"
            "3. 数学公式、符号和方程\n"
            "4. 手写内容\n"
            "5. 适合后续复习的关键点\n"
            "只描述图片中能看见的内容，不要凭空补充。"
        )
        return await llm.chat_with_image(
            system="你负责从图片中精确提取教育内容，并整理成可用于学习复习的文本。",
            messages=[{"role": "user", "content": prompt}],
            image_b64=b64,
            mime=mime,
        )
    except Exception as e:
        logger.error(f"Image parse error: {e}")
        return ""


# ── Code ──

def parse_code(file_path: str) -> str:
    """Read code file with language marker."""
    ext = Path(file_path).suffix.lower()
    lang = ext.lstrip(".")
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        return f"```{lang}\n{content}\n```"
    except Exception as e:
        logger.error(f"Code read error: {e}")
        return ""


# ── Text ──

def parse_text(file_path: str) -> str:
    """Read plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception as e:
        logger.error(f"Text read error: {e}")
        return ""


# ── Main dispatch ──

def parse_file(file_path: str, llm=None) -> dict:
    """Parse any supported file and return structured result.
    NOTE: For image files and scanned PDFs that need Vision, use parse_file_async() instead.

    Returns: {
        "ok": True/False,
        "text": "...",
        "type": "text_pdf" | "scanned_pdf" | "docx" | "pptx" | "image" | "code" | "text",
        "is_scanned": bool,
        "error": "..." | None
    }
    """
    ft = get_file_type(file_path)
    result = {"ok": True, "text": "", "type": ft, "is_scanned": False, "error": None}

    try:
        if ft == "text_pdf":
            text, is_scanned = parse_pdf(file_path)
            result["text"] = text
            result["is_scanned"] = is_scanned
            if is_scanned:
                result["type"] = "scanned_pdf"
        elif ft == "docx":
            result["text"] = parse_docx(file_path)
        elif ft == "pptx":
            result["text"] = parse_pptx(file_path)
        elif ft == "image":
            result["ok"] = False
            result["error"] = "Image files require parse_file_async() with an LLM"
        elif ft == "code":
            result["text"] = parse_code(file_path)
        elif ft == "text":
            result["text"] = parse_text(file_path)
        else:
            result["ok"] = False
            result["error"] = f"Unsupported file type: {ft}"
    except Exception as e:
        result["ok"] = False
        result["error"] = str(e)

    return result


async def parse_file_async(file_path: str, llm) -> dict:
    """Async version of parse_file — required for image/scanned PDF files that need Vision LLM.

    Returns: {
        "ok": True/False,
        "text": "...",
        "type": "text_pdf" | "scanned_pdf" | "docx" | "pptx" | "image" | "code" | "text",
        "is_scanned": bool,
        "error": "..." | None
    }
    """
    ft = get_file_type(file_path)
    result = {"ok": True, "text": "", "type": ft, "is_scanned": False, "error": None}

    try:
        if ft == "text_pdf":
            text, is_scanned = parse_pdf(file_path)
            result["text"] = text
            result["is_scanned"] = is_scanned
            if is_scanned and llm:
                result["type"] = "scanned_pdf"
                ocr_text = await parse_scanned_pdf(file_path, llm)
                if ocr_text and ocr_text.strip():
                    result["text"] = ocr_text
        elif ft == "docx":
            result["text"] = parse_docx(file_path)
        elif ft == "pptx":
            result["text"] = parse_pptx(file_path)
        elif ft == "image":
            if llm:
                result["text"] = await parse_image(file_path, llm)
            else:
                result["ok"] = False
                result["error"] = "No LLM available for image parsing"
        elif ft == "code":
            result["text"] = parse_code(file_path)
        elif ft == "text":
            result["text"] = parse_text(file_path)
        else:
            result["ok"] = False
            result["error"] = f"Unsupported file type: {ft}"
    except Exception as e:
        result["ok"] = False
        result["error"] = str(e)

    return result


def parse_first_pages(file_path: str) -> str:
    """Read first few pages/paragraphs/slides for classification."""
    ft = get_file_type(file_path)
    try:
        if ft == "text_pdf" or ft == "scanned_pdf":
            return parse_pdf_first_pages(file_path, pages=2)
        elif ft == "docx":
            return parse_docx_first_pages(file_path, paragraphs=15)
        elif ft == "pptx":
            return parse_pptx_first_pages(file_path, slides=3)
        elif ft == "image":
            return f"[Image] {Path(file_path).name}"
        elif ft in ("code", "text"):
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()[:2000]
        else:
            return ""
    except Exception as e:
        logger.error(f"First-pages error for {file_path}: {e}")
        return ""
